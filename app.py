from flask import Flask, request, jsonify, send_file, render_template
from flask_cors import CORS
import openai
import traceback
import os
import re
from PyPDF2 import PdfReader
from pptx import Presentation
from werkzeug.utils import secure_filename
from pptx.util import Inches, Pt
import time
import fitz  # Import PyMuPDF
import io
from PIL import Image

global file_path

app = Flask(__name__)
CORS(app)

openai.api_key = "Your api key here"


def split_text_by_separator(text, separator="$!$"):
    return text.split(separator)


def convert_pdf_to_images(pdf_path):
    doc = fitz.open(pdf_path)
    images = []
    for i in range(len(doc)):
        page = doc.load_page(i)
        pixmap = page.get_pixmap()
        temp_file_name = f"temp_image_{i}.png"
        pixmap.save(temp_file_name)  # Save the pixmap to a temporary file
        with open(temp_file_name, "rb") as image_file:
            image_bytes = io.BytesIO(image_file.read())  # Read the image file into a BytesIO buffer
        image = Image.open(image_bytes)  # Now open the image using PIL
        images.append(image)
        os.remove(temp_file_name)  # Delete the temporary file
    return images

def generate_pptx(processed_chunks, valid_images, output_filename):

    prs = Presentation()
    prs.slide_width = Inches(8.5)
    prs.slide_height = Inches(11)
    
    images = convert_pdf_to_images(file_path)  # Convert all PDF pages to images
    print(file_path)

    for i, processed_chunk in enumerate(processed_chunks):
        print(f"Processed chunk {i}:\n{processed_chunk}\n")  # Debug print

        # Updated regular expressions
        questions_with_answers = re.findall(r"Question:([\s\S]+?)Correct Answer:", processed_chunk, re.I)
        explanations = re.findall(r"Explanation:([\s\S]+?)(?=Question:|$)", processed_chunk, re.I)
        correct_answers = re.findall(r"Correct Answer:\s*([^\n]+)", processed_chunk, re.I)

        # Combine the correct answers with explanations
        explanations_with_correct_answers = []
        for idx, exp in enumerate(explanations):
            correct_answer = f"Correct Answer: {correct_answers[idx]}\n" if idx < len(correct_answers) else ""
            updated_exp = correct_answer + "Explanation: " + exp.strip()
            explanations = re.findall(r"Explanation:([\s\S]+)$", processed_chunk, re.I)
            explanations_with_correct_answers.append(updated_exp)

        print(f"Questions with Answers:\n{questions_with_answers}\n")  # Debug print
        print(f"Explanations with Correct Answers:\n{explanations_with_correct_answers}\n")  # Debug print

        # This loop should be inside the loop over processed_chunks
        for idx, (qa, e) in enumerate(zip(questions_with_answers, explanations_with_correct_answers)):
            # Create slides
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = f"Slide {i + 1}-{idx + 1} - Question"
            content.text = qa.strip()
            
            # Adjust the content position and size
            content.top = Inches(1.5)
            content.width = Inches(7.5)
            content.height = Inches(8)

            for paragraph in content.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(24)

            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = f"Slide {i + 1}-{idx + 1} - Answer"
            content.text = f"{e.strip()}"
            
            # Adjust the content position and size
            content.top = Inches(1.5)
            content.width = Inches(7.5)
            content.height = Inches(8)

            for paragraph in content.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(20)

        # Handle image addition
        image = valid_images[i]
        image_path = f'temp_page_{i}.png'
        image.save(image_path)
        slide_layout = prs.slide_layouts[5]  # Use the blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.add_picture(image_path, Inches(0), Inches(0), height=Inches(11))

    prs.save(output_filename)
    
    # Clean up temporary images
    output_dir = os.path.dirname(os.path.abspath(output_filename))
    for filename in os.listdir(output_dir):
        if filename.startswith('temp_page_') and filename.endswith('.png'):
            os.remove(os.path.join(output_dir, filename))



MAX_RETRIES = 5  # Max number of retries for API calls
DELAY_BETWEEN_RETRIES = 5  # Time delay in seconds between retries

def process_text(full_prompt):
    retries = 0
    selected_model = request.json["model"]
    while retries < MAX_RETRIES:
        try:
            system_message = {
                "role": "system",
                "content": f"You are Document GPT. Your job is to take the uploaded material and do with it what the user prompts you to."
            }
            user_message = {
                "role": "user",
                "content": full_prompt
            }

            response = openai.ChatCompletion.create(
                model=selected_model,  # Use the selected_model parameter here
                messages=[system_message, user_message],
                max_tokens=round(7500-len(full_prompt)/4),
                n=1,
                stop=None,
                temperature=0.5,
            )
            return response.choices[0].message['content'].strip()

        except Exception as e:
            retries += 1
            print(f"An error occurred: {e}. Retrying {retries}/{MAX_RETRIES}...")
            time.sleep(DELAY_BETWEEN_RETRIES)

@app.route('/process_transcript', methods=['POST'])
def process_transcript():
    try:
        transcript = request.json['transcript']
        prompt = request.json['gptPrompt']
        processing_mode = request.json['processingMode']
        page_content = request.json['pageContent']

        def split_text_by_pagecontent(text, page_content):
            separator = "$!$"
            chunks = []
            start_idx = 0
            
            while True:
                page_idx = text.find(page_content, start_idx)
                if page_idx == -1:
                    break
                
                separator_idx = text.find(separator, page_idx + len(page_content))
                if separator_idx != -1:
                    chunks.append(text[start_idx:separator_idx])
                    start_idx = separator_idx + len(separator)
                else:
                    chunks.append(text[start_idx:])
                    break

            return chunks

        if processing_mode == 'bySlide':
            # Split transcript by separator "$!$" if "bySlide" is selected
            transcript_chunks = split_text_by_separator(transcript)

            processed_chunks = []
            for chunk in transcript_chunks:
                full_prompt = f"\n\nPrompt: {prompt}\n\nUploaded Material:\n\n{chunk}"

                organized_chunk = process_text(full_prompt)
                processed_chunks.append(organized_chunk)

            organized_text = ' '.join(processed_chunks)
        elif processing_mode == 'wholeFile':
            transcript_chunks = split_text_by_pagecontent(transcript, page_content)

            processed_chunks = []
            for chunk in transcript_chunks:
                full_prompt = f"\n\nPrompt: {prompt}\n\nUploaded Material:\n\n{chunk}"
                organized_chunk = process_text(full_prompt)
                processed_chunks.append(organized_chunk)

            organized_text = ' '.join(processed_chunks)

        else:
            return jsonify({"error": "Invalid processing mode."}), 400

        result = {"processed_text": organized_text}
        return jsonify(result)
    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

@app.route('/process_practice_questions', methods=['POST'])
def process_practice_questions():
    try:
        transcript = request.json['transcript']
        transcript_chunks = split_text_by_separator(transcript)

        undesired_words = []

        
        valid_chunks = []
        images = convert_pdf_to_images(file_path)

        # Filtering chunks and corresponding images
        valid_images = []
        for idx, chunk in enumerate(transcript_chunks):
            if len(chunk.split()) < 5:
                continue

            if any(word in chunk.lower() for word in undesired_words):
                continue

            valid_chunks.append(chunk)
            valid_images.append(images[idx])

        processed_chunks = []
        for chunk in valid_chunks:
            full_prompt = f'''Prompt: "Generate a USMLE style medical school exam question based on the slide details provided:

Clinically Relevant and Direct Content Extraction: Focus on crafting a question around the main theme or core information presented in the slide. Develop a clinically relevant scenario and ensure that the question, its answer choices, and its explanation are derived primarily from central concepts of the slide, focusing on mutations, imaging, histology, pathogenesis, anatomy/physiology, treatments, pharmacology, or clinical presentation.

Deep Contextual Understanding: The question should resonate with the general theme and overarching message of the slide content. Both the question and answer choices should be based on the slide's primary material, ensuring that the correct answer is challenging but not overly obscure. Avoid explicitly providing a diagnosis, rather focusing on clinical presentation within the question stem.

Second-Order Thinking: The question should stimulate students to connect a clinical presentation to its underlying pathology, mechanism, or treatment. This connection should be based on major concepts in the slide, without leaning too heavily on minutiae.

Subject Range: Depending on the slide, frame questions around broad topics such as mutations, imaging, histology, pathogenesis, anatomy/physiology, treatments, pharmacology, or clinical presentation.

Consistency & Accuracy: Align the question, answer, and slide content cohesively. There should be only one unmistakably correct answer.

Balanced Answer Choices: Maintain similar detail levels for all answer choices to prevent the correct answer from standing out due to detail disparity. Ensure answer choices are factually accurate and reflect the content of the slide. 

Misconception-based Distractors: Integrate incorrect options that reflect common misconceptions, ensuring they are derived from or closely related to the slide's main content and ensuring they are not repetitive. Answer choices should be singular rather than asking to identify two components. 

Caution on Comparative Questions: If the question involves comparing or prioritizing different factors (e.g., highest risk, most significant factor, next best step), ensure that the slide content provides explicit information to support the comparison or hierarchy. Avoid making unsupported inferences or assumptions that are not directly based on the slide content.

Answer Choices: Refrain from using aggregated choices like "all of the above." Each answer choice should be distinct. Avoid explicitly asking for a diagnosis. Answer choices must require second order thought process while still relying explictly on factual information provided within the slide

Format:

Question: [Question]
Answer Choices: [A-E]
Correct Answer: [correct answer]
Explanation: A well-rounded paragraph that explains the rationale behind the correct answer and provides succinct reasons for the incorrectness of each of the other choices."\nUploaded Material: {chunk}'''

            organized_chunk = process_text(full_prompt)
            processed_chunks.append(organized_chunk)

        original_filename = os.path.basename(file_path).rsplit('.', 1)[0]
        output_filename = f"{original_filename} practice questions.pptx"
        generate_pptx(processed_chunks, valid_images, output_filename)
        return "File Saved"

    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500


UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'pdf', 'pptx'}

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# Modified function to include separator "$!$"
def extract_text_from_pdf(file_path):
    pages_text = []
    with open(file_path, 'rb') as file:
        pdf = PdfReader(file)
        n_pages = len(pdf.pages)
        for page_number in range(n_pages):
            page = pdf.pages[page_number]
            pages_text.append(page.extract_text())
    text = "$!$".join(pages_text)  # separator will not be added at the end
    return text

# Modified function to include separator "$!$"

def extract_text_from_pptx(file_path):
    text = ""
    prs = Presentation(file_path)

    for slide in prs.slides:
        slide_text = ""
        notes_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text += run.text + ' '

        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            for shape in notes_slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            notes_text += run.text + ' '

        # Combine slide text and notes text for the current slide, and add the separator
        text += slide_text + ' ' + notes_text + "$!$"

    return text

@app.route('/upload_file', methods=['POST'])
def upload_file():
    global file_path  # Indicate that we're using the global file_path variable

    # Clear everything in the uploads folder
    folder = app.config['UPLOAD_FOLDER']
    for filename in os.listdir(folder):
        file_path = os.path.join(folder, filename)
        try:
            if os.path.isfile(file_path) or os.path.islink(file_path):
                os.unlink(file_path)
            elif os.path.isdir(file_path):
                shutil.rmtree(file_path)
        except Exception as e:
            print(f'Failed to delete {file_path}. Reason: {e}')

    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No file selected"}), 400
    if not allowed_file(file.filename):
        return jsonify({"error": "File type not allowed"}), 400
    filename = secure_filename(file.filename)
    file.save(os.path.join(folder, filename))

    file_path = os.path.join(folder, filename)  # Assign to global file_path
    file_extension = filename.rsplit('.', 1)[1].lower()

    if file_extension == 'pdf':
        text = extract_text_from_pdf(file_path)
    elif file_extension == 'pptx':
        text = extract_text_from_pptx(file_path)
    else:
        os.remove(file_path)
        return jsonify({"error": "Unsupported file type"}), 400

    return jsonify({"text": text})

@app.route('/')
def index():
    return render_template('index.html')

if __name__ == "__main__":
    app.run(debug=True)
