import os
import comtypes.client
import PyPDF2
import random
from io import BytesIO
import fitz
from PIL import Image
from pptx.util import Inches
from pptx import Presentation
import io
import time

def ppt_to_pdf(input_path, output_path):
    """
    Convert a PowerPoint file to PDF using COM automation with Microsoft PowerPoint.
    
    :param input_path: str, path to the PowerPoint file.
    :param output_path: str, path to save the output PDF file.
    """
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    if not os.path.isabs(input_path):
        input_path = os.path.abspath(input_path)

    if not os.path.isabs(output_path):
        output_path = os.path.abspath(output_path)

    deck = powerpoint.Presentations.Open(input_path)
    deck.SaveAs(output_path, 32)  # 32 corresponds to the format for PDFs in PowerPoint's FileFormat enumeration
    deck.Close()
    powerpoint.Quit()
    
def split_pdf_into_units(pdf_path):
    """Split a PDF into units of three pages."""
    with open(pdf_path, 'rb') as pdf_file:
        reader = PyPDF2.PdfReader(pdf_file)
        total_pages = len(reader.pages)
        units = []

        for i in range(0, total_pages, 3):
            output = BytesIO()
            unit = PyPDF2.PdfWriter()
            for j in range(3):
                if (i + j) < total_pages:
                    unit.add_page(reader.pages[i + j])
            unit.write(output)
            output.seek(0)
            units.append(output)
    return units

def merge_units_into_pdf(units, output_path):
    """Merge a list of 3-page units into a single PDF."""
    with open(output_path, 'wb') as output_file:
        merger = PyPDF2.PdfMerger()
        for unit_stream in units:
            merger.append(unit_stream)
        merger.write(output_file)

def process_pdfs(pdf_paths, output_path):
    """Process multiple PDFs, split them, randomize the units, and merge them."""
    all_units = []
    for pdf_path in pdf_paths:
        all_units.extend(split_pdf_into_units(pdf_path))

    random.shuffle(all_units)
    merge_units_into_pdf(all_units, output_path)
    
def convert_pdf_to_images(pdf_path):
    doc = fitz.open(pdf_path)
    images = []
    for i in range(len(doc)):
        page = doc.load_page(i)
        pixmap = page.get_pixmap()
        temp_file_name = f"temp_image_{i}.png"
        pixmap.save(temp_file_name)  # Save the pixmap to a temporary file
        with open(temp_file_name, "rb") as image_file:
            image_bytes = io.BytesIO(image_file.read())  # Read the image file into a BytesIO buffer
        image = Image.open(image_bytes)  # Now open the image using PIL
        images.append(image)
        os.remove(temp_file_name)  # Delete the temporary file
    return images

def images_to_pptx(images, pptx_path):
    prs = Presentation()

    # Define slide width and height (8x11 inches)
    prs.slide_width = Inches(8.5)
    prs.slide_height = Inches(11)

    for i, image in enumerate(images):
        image_path = f'temp_page_{i}.png'
        image.save(image_path)

        slide_layout = prs.slide_layouts[5]  # use the blank slide layout
        slide = prs.slides.add_slide(slide_layout)

        left = Inches(0)
        top = Inches(0)
        height = Inches(11)  # setting height to maintain aspect ratio
        pic = slide.shapes.add_picture(image_path, left, top, height=height)
        os.remove(image_path)  # Delete the temporary image file

    prs.save(pptx_path)
    
def is_powerpoint(file_name):
    """Check if the given file name is a PowerPoint file."""
    return file_name.endswith(('.ppt', '.pptx'))

def num_slides_in_ppt(file_path):
    """Return the number of slides in a PowerPoint file."""
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    deck = powerpoint.Presentations.Open(file_path)
    num_slides = deck.Slides.Count
    powerpoint.Quit()
    return num_slides

def process_folder(folder_path):
    all_units = []  # This will store chunks from all PowerPoints
    
    # Convert each PowerPoint in the folder to a PDF and chunk it
    for file_name in os.listdir(folder_path):
        file_path = os.path.join(folder_path, file_name)
        
        # Check if it's a PowerPoint file
        if is_powerpoint(file_name):
            # Check if the number of slides is divisible by three
            if num_slides_in_ppt(file_path) % 3 == 0:
                pdf_path = os.path.join(folder_path, f"temp_{file_name.replace('.pptx', '').replace('.ppt', '')}.pdf")
                
                # Convert to PDF
                ppt_to_pdf(file_path, pdf_path)
                
                # Chunk the PDF into units and add to all_units
                all_units.extend(split_pdf_into_units(pdf_path))
                
                # Delete the temporary PDF
                os.remove(pdf_path)
    
    # Randomize all chunks
    random.shuffle(all_units)
    
    # Merge the randomized chunks into a single PDF
    merged_pdf_path = os.path.join(folder_path, 'random_temp.pdf')
    merge_units_into_pdf(all_units, merged_pdf_path)
    
    # Convert the merged PDF back to PowerPoint
    images = convert_pdf_to_images(merged_pdf_path)
    pptx_path = os.path.join(folder_path, 'Random.pptx')
    images_to_pptx(images, pptx_path)
        # Add a delay to ensure all processes using the file have released it
    time.sleep(2)
    # Delete the temporary merged PDF
    os.remove(merged_pdf_path)

if __name__ == "__main__":
    # Get the directory where the script is located
    folder_path = os.path.dirname(os.path.abspath(__file__))
    process_folder(folder_path)
